import os
import requests
import uuid
import PyPDF2
from pptx import Presentation
import json
import requests
import time
import re
from gtts import gTTS
from pdf2image import convert_from_path
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips, VideoFileClip
from multiprocessing import Pool

# Function to extract text from PDF or PPTX
def extract_text_by_file_type(file_path, output_directory):
    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()

    if file_extension == '.pdf':
        extract_pdf_text(file_path, output_directory)
    elif file_extension == '.pptx':
        extract_ppt_text(file_path, output_directory)
    else:
        raise ValueError('Unsupported file format: {}'.format(file_extension))

def extract_pdf_text(file_path, output_directory):
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        num_pages = len(reader.pages)
        for page_number in range(num_pages):
            page = reader.pages[page_number]
            text = page.extract_text()
            output_file_path = os.path.join(output_directory, f'extracted_text_{page_number + 1}.txt')
            with open(output_file_path, 'w', encoding='utf-8') as f:
                json.dump({page_number + 1: text}, f, ensure_ascii=False, indent=4)

def extract_ppt_text(file_path, output_directory):
    prs = Presentation(file_path)
    for idx, slide in enumerate(prs.slides):
        slide_text = ''
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + '\n'
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text += run.text
                    slide_text += '\n'
        output_file_path = os.path.join(output_directory, f'extracted_text_{idx + 1}.txt')
        with open(output_file_path, 'w', encoding='utf-8') as f:
            json.dump({idx + 1: slide_text.strip()}, f, ensure_ascii=False, indent=4)

# Function to generate scripts using the GPT API
def generate_scripts(input_directory, output_directory, selected_style, selected_token_limit):
    API_KEY = "3883181699334026977e55d8e18b1f4d"  # Replace with your actual API key
    headers = {
        "Content-Type": "application/json",
        "api-key": API_KEY,
    }
    ENDPOINT = "https://fyp-ik2402-eastus-standards0.openai.azure.com/openai/deployments/gpt-4o-mini-globalstandard/chat/completions?api-version=2024-02-15-preview"  # Replace with your API endpoint

    def send_text_to_api(page_text, page_number, style, token_limit):
        prompt = f"You are an AI assistant that helps teachers generate explanatory scripts in a {style} style. Please ensure the script is strictly less than {token_limit} tokens."
        payload = {
            "messages": [
                {"role": "system", "content": prompt},
                {"role": "user", "content": page_text}
            ],
            "temperature": 0.7,
            "top_p": 0.95,
            "max_tokens": token_limit,
        }
        for attempt in range(3):  # Retry up to 3 times
            try:
                response = requests.post(ENDPOINT, headers=headers, json=payload, timeout=10)
                response.raise_for_status()
                generated_content = response.json()['choices'][0]['message']['content']
                return generated_content
            except requests.exceptions.RequestException as e:
                print(f"Error for page {page_number} (attempt {attempt + 1}): {e}")
                time.sleep(5)  # Wait 5 seconds before retrying
            except KeyError as e:
                print(f"KeyError for page {page_number} (attempt {attempt + 1}): {e}")
                time.sleep(5)
        return None  # Return None if all attempts fail

    for file_name in os.listdir(input_directory):
        if file_name.endswith('.txt'):
            page_number = int(file_name.split('_')[-1].replace('.txt', ''))
            input_file_path = os.path.join(input_directory, file_name)
            with open(input_file_path, 'r', encoding='utf-8') as f:
                content = f.read()
                data = json.loads(content)
                page_text = data.get(str(page_number), "")
            if page_text:
                generated_content = send_text_to_api(page_text, page_number, selected_style, selected_token_limit)
                if generated_content:
                    # 去除生成内容中的换行符
                    generated_content = generated_content.replace('\n', ' ').replace('\r', '')
                    # 调用 extract_core_content 过滤生成的内容
                    filtered_content = extract_core_content(generated_content)
                   
                    output_file_path = os.path.join(output_directory, f'generated_content_{page_number}.txt')
                    with open(output_file_path, 'w', encoding='utf-8') as output_file:
                        output_file.write(filtered_content)
                else:
                    print(f"Failed to generate content for page {page_number}")

# Function to filter content


'''
这个函数过滤的程度有些大
有一定几率会导致生成的generated script为空
因为要降低过滤的程度，防止上述的情况出现

def extract_core_content(text):
    # Exclude patterns to filter out unwanted lines
    exclude_patterns = [
        r'^Page\s+\d+',  # Lines starting with 'Page [number]'
        r'^---+$',       # Lines that are just '---'
        r'^\s*$',        # Empty lines
        r'^(Sure|Certainly|Let me|Here is|Feel free|This script|Hope this helps|Let me know if|Hello|That\'s|In summary|To sum up|If you have any questions|Thank you|So,|Now,|Remember,|Also,|Overall,|Lastly,|Finally,|First,|Second,|Third|Additionally|Moreover|Therefore|Thus|As you can see|I hope this|Don’t hesitate|Keep in mind|Feel free).*',  # Unwanted opening phrases
        r'.*(If you have any questions|feel free to ask|need more examples).*',  # Unwanted phrases within lines
        r'^#+.*',        # Lines starting with '#' (markdown headings)
        r'^\*\*.*\*\*$', # Lines that are entirely bolded (e.g., '**Title**')
    ]
    exclude_regex = [re.compile(p, re.IGNORECASE) for p in exclude_patterns]
    
    # Split text into lines and filter out unwanted lines
    lines = text.splitlines()
    filtered_lines = []
    for line in lines:
        line_stripped = line.strip()
        if any(regex.match(line_stripped) for regex in exclude_regex):
            continue
        filtered_lines.append(line)
    
    # Rejoin the filtered lines
    filtered_text = '\n'.join(filtered_lines)

    # Remove markdown formatting
    filtered_text = re.sub(r'\*\*(.*?)\*\*', r'\1', filtered_text)  # Remove bold (**text**)
    filtered_text = re.sub(r'\*(.*?)\*', r'\1', filtered_text)       # Remove italics (*text*)
    filtered_text = re.sub(r'^#+\s*(.*)', r'\1', filtered_text, flags=re.MULTILINE)  # Remove headings (# Text)
    filtered_text = re.sub(r'^---+$', '', filtered_text, flags=re.MULTILINE)  # Remove horizontal rules
    filtered_text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', filtered_text)  # Remove links [text](url)
    
    # Strip extra whitespace
    filtered_text = '\n'.join(line.strip() for line in filtered_text.strip().splitlines() if line.strip())
    
    return filtered_text

'''


import os
import requests

def convert_text_to_speech_with_clone(input_dir, output_dir, ref_audio, ref_text):
    for file_name in os.listdir(input_dir):
        if file_name.endswith('.txt'):
            input_path = os.path.join(input_dir, file_name)
            with open(input_path, 'r', encoding='utf-8') as f:
                text_content = f.read()
            
            file_number = file_name.split('_')[-1].replace('.txt', '')
            output_file_path = os.path.join(output_dir, f'clone_speech_{file_number}.mp3')

            # 构建请求数据
            data = {
                "text": text_content,
                "reference_audio":ref_audio,
                "reference_text": ref_text  # 同语言克隆必须参数
            }
            response=requests.post(f'http://127.0.0.1:9233/clone_eq',data=data,timeout=3600)
                
            if response.status_code == 200:
                with open(output_file_path, 'wb') as f:
                        f.write(response.content)
                print(f"Cloned audio saved to {output_file_path}")
            else:
                print(f"API Error [{response.status_code}]: {response.text}")
                print(response.json())


            
def extract_core_content(text):
    # 更精准的排除模式
    exclude_patterns = [
        r'^Page\s+\d+\s*$',  # 严格匹配 "Page X" 行
        r'^-{3,}$',          # 匹配至少3个连字符的分隔线
        r'^\s*$',            # 空行
        # 更精准的短语匹配（仅匹配以这些短语开头的行）
        r'^(Sure,|Certainly!|Let me explain|Here is (a|the)|Feel free to|This script|Hope this helps|Hello!|That\'s |In summary:|To sum up,|If you have any\b|Thank you|So, |Now, |Remember: |Also, |Overall,|Lastly, |Finally, |First,|Second,|Third:|Additionally, |Moreover, |Therefore, |Thus, |As you can see, |I hope this\b|Don’t hesitate\b)',
        # 不再排除Markdown格式行（后续用替换处理）
    ]
    
    exclude_regex = [re.compile(p, re.IGNORECASE) for p in exclude_patterns]
    
    # 处理流程优化
    lines = text.splitlines()
    filtered_lines = []
    
    for line in lines:
        line_stripped = line.strip()
        
        # 跳过完全匹配排除模式的行
        if any(regex.match(line_stripped) for regex in exclude_regex):
            continue
            
        filtered_lines.append(line)
    
    filtered_text = '\n'.join(filtered_lines)
    
    # 保留内容但移除格式（不再暴力删除整行）
    replacements = [
        (r'\*\*(.*?)\*\*', r'\1'),   # 保留加粗内容
        (r'\*(.*?)\*', r'\1'),       # 保留斜体内容
        (r'^#+\s*(.*)', r'\1'),      # 保留标题内容
        (r'\[(.*?)\]\(.*?\)', r'\1') # 保留链接文字
    ]
    
    for pattern, repl in replacements:
        filtered_text = re.sub(pattern, repl, filtered_text, flags=re.MULTILINE)
    
    # 清理多余空行并保留段落结构
    filtered_text = re.sub(r'\n\s*\n', '\n\n', filtered_text.strip())
    
    return filtered_text



# Function to convert text to speech

import requests

def convert_text_to_speech_without_clone(input_directory, output_directory,voice_type):
    for file_name in os.listdir(input_directory):
        if file_name.endswith('.txt'):
            input_file_path = os.path.join(input_directory, file_name)
            with open(input_file_path, 'r', encoding='utf-8') as file:
                file_content = file.read()
            file_number = file_name.split('_')[-1].replace('.txt', '')
            output_file_path = os.path.join(output_directory, f'speech_{file_number}.mp3')  # 保持与原来相同的文件名和扩展名
        
            # 准备发送到TTS API的数据
            data = {
                "text": file_content,
                "role":voice_type # 根据你的API需求，这里可能需要调整
            }

            # 发送POST请求到TTS API
            try:
                response = requests.post('http://127.0.0.1:9233/tts', data=data, timeout=3600)
                if response.status_code != 200:
                    # 出错了
                    print(f"Error: {response.json()}")
                else:
                    # 返回的wav数据流，可直接保存
                    with open(output_file_path, 'wb') as f:
                        f.write(response.content)
                    print(f"Speech file saved as {output_file_path}")
            except requests.exceptions.RequestException as e:
                print(f"Request failed: {e}")

# Function to convert PDF to images

def pdf_to_images(pdf_path, output_dir):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    images = convert_from_path(pdf_path, dpi=150, thread_count=4)
    num_pages = len(images)  # 统计 PDF 总页数
    for i, page in enumerate(images):
        image_path = os.path.join(output_dir, f'page_{i+1}.png')
        page.save(image_path, 'PNG')
    
    return num_pages

#新增函数
def get_num_scripts_in_dir(generated_script_dir):
    """
    统计 generated_script_dir 下有多少脚本文件（.txt)。
    """
    return len([f for f in os.listdir(generated_script_dir) if f.endswith(".txt")])

# Function to create the presentation video
def create_presentation_video(image_dir, audio_dir, output_path):
    image_files = [os.path.join(image_dir, f) for f in os.listdir(image_dir) if f.endswith('.png')]
    audio_files = [os.path.join(audio_dir, f) for f in os.listdir(audio_dir) if f.endswith('.mp3')]

    def extract_number(filename):
        match = re.search(r'(\d+)', filename)
        return int(match.group(1)) if match else -1

    image_files.sort(key=lambda x: extract_number(os.path.basename(x)))
    audio_files.sort(key=lambda x: extract_number(os.path.basename(x)))

    if len(image_files) != len(audio_files):
        print("Mismatch between image and audio files.")
        return

    temp_video_dir = os.path.join(os.path.dirname(output_path), "temp_videos")
    if not os.path.exists(temp_video_dir):
        os.makedirs(temp_video_dir)

    args_list = [(img, aud, temp_video_dir) for img, aud in zip(image_files, audio_files)]

    with Pool() as pool:
        video_paths = pool.map(create_slide_video, args_list)

    video_paths = [vp for vp in video_paths if vp is not None]
    clips = [VideoFileClip(vp) for vp in video_paths]
    final_video = concatenate_videoclips(clips, method="chain")
    final_video.write_videofile(output_path, fps=24, preset='ultrafast', threads=4, codec='libx264', audio_codec='libmp3lame')

def create_slide_video(args):
    image_path, audio_path, output_dir = args
    try:
        audio_clip = AudioFileClip(audio_path)
        duration = audio_clip.duration
        image_clip = ImageClip(image_path).set_duration(duration).resize(height=720).set_audio(audio_clip)
        base_name = os.path.basename(image_path)
        output_filename = os.path.splitext(base_name)[0] + '_slide.mp4'
        output_path = os.path.join(output_dir, output_filename)
        image_clip.write_videofile(
            output_path,
            fps=24,
            preset='ultrafast',
            threads=4,
            codec='libx264',
            audio_codec='libmp3lame',
            logger='bar'  # Enable logging
        )
        return output_path
    except Exception as e:
        print(f"Error in create_slide_video for {image_path} and {audio_path}: {e}")
        return None
    



# utils.py
import hashlib
import hmac
import base64
import json
import time
import requests
from requests_toolbelt.multipart.encoder import MultipartEncoder

class AIPPT:
    def __init__(self, app_id, api_secret, text=None, 
                 template_id='202407171E27C9D', file_name=None, 
                 file_content=None, file_url=None, language='cn'):
        self.APPid = app_id
        self.APISecret = api_secret
        self.text = text
        self.templateId = template_id
        self.file_name = file_name
        self.file_content = file_content
        self.file_url = file_url
        self.language = language  # 新增语言参数
        self.header = {}

    def get_signature(self, ts):
        auth = hashlib.md5(f"{self.APPid}{ts}".encode()).hexdigest()
        return base64.b64encode(
            hmac.new(self.APISecret.encode(), auth.encode(), hashlib.sha1).digest()
        ).decode()

    def create_task(self):
        url = 'https://zwapi.xfyun.cn/api/ppt/v2/create'
        timestamp = int(time.time())
        signature = self.get_signature(timestamp)
        
        # 构建表单字段
        fields = {
            "templateId": self.templateId,
            "language": self.language,
            "author": "AutoGenerated",
            "isCardNote": "True",
            "search": "False",
            "isFigure": "True",
            "aiImage": "normal"
        }
        

        #这里的逻辑是：
        #只传入文件，用文件内容
        #只写文字，用文字内容
        #如果又上传了文件，又有文字内容，两者都用
        if self.file_content and self.file_name:
            fields["file"] = (self.file_name, self.file_content, 'application/octet-stream')
            fields["fileName"] = self.file_name
            if self.text:  # 附加文本作为补充
                fields["query"] = self.text
        else:
            fields["query"] = self.text

        form_data = MultipartEncoder(fields=fields)
        
        headers = {
            "appId": self.APPid,
            "timestamp": str(timestamp),
            "signature": signature,
            "Content-Type": form_data.content_type
        }

        response = requests.post(url, data=form_data, headers=headers)
        response_data = json.loads(response.text)
        
        if response_data.get('code') == 0:
            return response_data['data']['sid']
        raise Exception(response_data.get('message', 'PPT creation failed'))

    def get_result(self, task_id):
        while True:
            progress = json.loads(self.get_process(task_id))
            if all([
                progress['data']['pptStatus'] == 'done',
                progress['data']['aiImageStatus'] == 'done',
                progress['data']['cardNoteStatus'] == 'done'
            ]):
                return progress['data']['pptUrl']
            time.sleep(3)

    def get_process(self, sid):
        url = f"https://zwapi.xfyun.cn/api/ppt/v2/progress?sid={sid}"
        timestamp = int(time.time())
        signature = self.get_signature(timestamp)
        
        headers = {
            "appId": self.APPid,
            "timestamp": str(timestamp),
            "signature": signature,
        }
        
        response = requests.get(url, headers=headers)
        return response.text